from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

NAVY = RGBColor(0x0A, 0x32, 0x56)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ICE = RGBColor(0xE6, 0xF2, 0xFF)
TEXT = RGBColor(0x00, 0x33, 0x59)

prs = Presentation('interim.pptx')
slides = list(prs.slides)

def find_ph(slide, idx=None, ph_type=None):
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        pf = shape.placeholder_format
        if idx is not None and pf.idx != idx:
            continue
        if ph_type is not None and pf.type != ph_type:
            continue
        return shape
    return None

def set_text(shape, text):
    tf = shape.text_frame
    # keep first run's formatting, overwrite text
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]:
            r.text = ''
    else:
        p.text = text

def style_table(table, col_widths_in, row_heights_in, headers, rows, font_size=10, header_font_size=11):
    for i, w in enumerate(col_widths_in):
        table.columns[i].width = Inches(w)
    for i, h in enumerate(row_heights_in):
        table.rows[i].height = Inches(h)

    # header row
    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = htext
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(6); cell.margin_right = Pt(6)
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.runs[0]
        run.font.size = Pt(header_font_size)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = 'Arial'

    # body rows
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ICE if r % 2 == 0 else WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Pt(6); cell.margin_right = Pt(6)
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            run = para.runs[0]
            run.font.size = Pt(font_size)
            run.font.bold = False
            run.font.color.rgb = TEXT
            run.font.name = 'Arial'

# ---------------------------------------------------------------------------
# Slide index 4 (0-based) = slide32.xml = "7-Week Plan to Go-Live"
# ---------------------------------------------------------------------------
s = slides[3]
title = find_ph(s, idx=None, ph_type=None)
for shape in s.shapes:
    if shape.is_placeholder and shape.placeholder_format.type is not None and str(shape.placeholder_format.type) == 'TITLE (13)':
        pass

# locate by idx directly (title has no idx attr in pptx model -> idx 0 conventionally not guaranteed; use type name match instead)
def find_by_type_name(slide, type_name_contains):
    for shape in slide.shapes:
        if shape.is_placeholder and type_name_contains in str(shape.placeholder_format.type):
            return shape
    return None

title_shape = find_by_type_name(s, 'TITLE')
set_text(title_shape, '7-Week Plan to Go-Live')

subheader_shape = find_ph(s, idx=15)
set_text(subheader_shape, 'Named owners, weekly milestones, dependencies')

content_shape = find_ph(s, idx=1)
left, top, width, height = content_shape.left, content_shape.top, content_shape.width, content_shape.height
content_shape._element.getparent().remove(content_shape._element)

headers = ['Week', 'Milestone', 'Owner', 'Dependency']
rows = [
    ['5', 'Snowflake read access granted; SSO working session held (SAML, scoping, mTLS); tablet shared-device model confirmed', 'Raj / FDE + Lisa + Marcus / Tom + Diana', 'Lisa back from leave'],
    ['6', 'Steering committee sign-off on integration approach; Snowflake schema/freshness validated; mTLS cert request filed', 'Sarah / RELEX Data Eng / Marcus', 'This meeting'],
    ['7', 'SSO (SAML) implemented end-to-end; UPC dedup rule built into ingestion pipeline', 'FDE + RELEX Data Eng', 'Week 6 sign-off'],
    ['8', 'Per-store order file emitter built; ingestion pipeline live in staging', 'RELEX Eng', 'Snowflake access'],
    ['9', 'First end-to-end test order: RELEX proposal → SFTP file → SAP acceptance, single pilot store', 'FDE + Raj', 'Order emitter built'],
    ['10', 'Expand test orders to all 25 stores; tablet app + SSO live for store managers', 'FDE + Tom', 'Week 9 test passes'],
    ['11', 'Parallel run — managers review RELEX proposals alongside manual ordering, no cutover yet', 'Tom + Diana', 'Week 10 stable'],
    ['12', 'Go-Live: RELEX orders flow live to SAP for 25 pilot stores', 'FDE + Sarah (go/no-go call)', 'Week 11 clean'],
]
gtable = s.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
table = gtable.table
col_widths = [0.8, 5.4, 3.0, 2.85]
row_h = height / Emu(Inches(1)) / len(table.rows)
style_table(table, col_widths, [row_h] * len(table.rows), headers, rows, font_size=9.5, header_font_size=11)

# ---------------------------------------------------------------------------
# Slide index 5 (0-based) = slide33.xml = "What We Need From Northbridge"
# ---------------------------------------------------------------------------
s = slides[4]
title_shape = find_by_type_name(s, 'TITLE')
set_text(title_shape, 'What We Need From Northbridge')

subheader_shape = find_ph(s, idx=15)
set_text(subheader_shape, 'Specific asks, owners, and dates — not open-ended risks')

content_shape = find_ph(s, idx=1)
left, top, width, height = content_shape.left, content_shape.top, content_shape.width, content_shape.height
content_shape._element.getparent().remove(content_shape._element)

headers = ['Ask', 'From', 'By']
rows = [
    ['Read-only Snowflake role for RELEX', 'Raj', 'Week 5'],
    ['90-min working session to close SSO (SAML, per-store scoping, mTLS)', 'Lisa + Marcus', 'Week 5'],
    ['Confirm whether Azure AD can source a store_id claim from Workday', 'Marcus', 'Week 5'],
    ['mTLS client certificate request filed (4–6 week lead time — start now)', "Lisa's security team", 'Requested Week 5, needed Week 10'],
    ['Decision: accept ORDER_TYPE = REGULAR tagging for pilot orders', 'Raj + Marcus', 'Week 6'],
    ['Confirm tablet shared-device / re-auth model', 'Tom + Diana', 'Week 5'],
    ['Sign-off on Snowflake-read + SFTP-order integration approach', 'Sarah', 'Today'],
]
gtable = s.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
table = gtable.table
col_widths = [7.3, 2.7, 2.05]
row_h = height / Emu(Inches(1)) / len(table.rows)
style_table(table, col_widths, [row_h] * len(table.rows), headers, rows, font_size=10.5, header_font_size=12)

prs.save('final.pptx')
print('Saved final.pptx')
